"""Turn a downloaded file into plain text, or into nothing at all.

Every reader here is pure-Python and adds no system libraries to the image.

Returning "" is a first-class outcome, not a failure. A scanned PDF with no text
layer, a legacy .doc, and a DocuSign form all have no readable body, and the
caller must treat that as "cannot read" rather than "the content is empty" —
that distinction is the difference between reporting a file and deleting a
document, which is the most destructive thing this job could do.
"""

from __future__ import annotations

import io
import logging
import re

log = logging.getLogger(__name__)

# Caps for the SCRAPE job, where extraction is a change-detection gate over ~277
# files every run and reading a 118-page policy in full buys nothing.
#
# They are the wrong caps for INGESTION, and the failure is silent: a truncated
# read returns clean text with no error, so a document looks complete and a
# "does the KB cover this?" check wrongly clears. Measured 2026-08-07 -- MSU's
# Procurement Policies is 118 pages / 245,608 chars and came back 36% captured
# (SECTIONS VI-XI and Appendix A missing); the Prohibited Conduct Procedures is
# 89 pages and came back at 40. Callers that are LOADING a file into the KB must
# pass max_pages/max_chars explicitly; the loader scripts pass None for both.
MAX_CHARS = 200_000
MAX_PDF_PAGES = 40

_FILE_EXT = (".pdf", ".doc", ".docx", ".xls", ".xlsx", ".ppt", ".pptx")
# A DocuSign PowerForm or a Google form is a live web page, not a document. It
# has no downloadable body, and its bytes differ on every request, so it can be
# linked but never hashed.
_FORM_HOSTS = ("docusign.net", "forms.gle", "docs.google.com/forms")


def kind_for_url(url: str) -> str:
    """file | form | page — what kind of destination this URL is."""
    low = (url or "").lower()
    if any(h in low for h in _FORM_HOSTS):
        return "form"
    path = low.split("?")[0].split("#")[0]
    if path.endswith(_FILE_EXT):
        return "file"
    return "page"


def _truncate(text: str, max_chars: int | None = MAX_CHARS) -> str:
    return text if max_chars is None else text[:max_chars]


def _clean(text: str, max_chars: int | None = MAX_CHARS) -> str:
    return _truncate(re.sub(r"\n{3,}", "\n\n", text).strip(), max_chars)


def _pdf(raw: bytes, max_pages: int | None = MAX_PDF_PAGES) -> str:
    import pdfplumber

    out = []
    with pdfplumber.open(io.BytesIO(raw)) as pdf:
        pages = pdf.pages if max_pages is None else pdf.pages[:max_pages]
        for page in pages:
            out.append(page.extract_text() or "")
    return "\n\n".join(out)


def _docx(raw: bytes) -> str:
    import docx

    d = docx.Document(io.BytesIO(raw))
    parts = [p.text for p in d.paragraphs]
    for table in d.tables:
        for row in table.rows:
            parts.append(" | ".join(c.text.strip() for c in row.cells))
    return "\n".join(parts)


def _pptx(raw: bytes) -> str:
    from pptx import Presentation

    parts = []
    for slide in Presentation(io.BytesIO(raw)).slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                parts.append(shape.text_frame.text)
    return "\n".join(parts)


def _xlsx(raw: bytes) -> str:
    import openpyxl

    wb = openpyxl.load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
    parts = []
    for sheet in wb.worksheets:
        parts.append(f"# {sheet.title}")
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


_READERS = {
    "pdf": _pdf,
    "docx": _docx,
    "pptx": _pptx,
    "xlsx": _xlsx,
    "xls": _xlsx,
}


def _format_of(content_type: str, url: str) -> str:
    """Prefer the served Content-Type; fall back to the URL's extension."""
    ct = (content_type or "").lower()
    if "pdf" in ct:
        return "pdf"
    if "wordprocessingml" in ct:
        return "docx"
    if "presentationml" in ct:
        return "pptx"
    if "spreadsheetml" in ct or "ms-excel" in ct:
        return "xlsx"
    path = (url or "").lower().split("?")[0].split("#")[0]
    for ext in ("pdf", "docx", "pptx", "xlsx", "xls"):
        if path.endswith("." + ext):
            return ext
    return ""


def extract_text(raw: bytes, content_type: str = "", url: str = "",
                 max_pages: int | None = MAX_PDF_PAGES,
                 max_chars: int | None = MAX_CHARS) -> str:
    """Plain text, or "" when this file has no readable body. Never raises.

    Pass max_pages=None, max_chars=None when INGESTING a file into the KB. The
    defaults are the scrape job's change-detection caps and will silently return
    a partial document -- see the note on MAX_CHARS.
    """
    if not raw:
        return ""
    reader = _READERS.get(_format_of(content_type, url))
    if reader is None:
        return ""
    try:
        if reader is _pdf:
            return _clean(_pdf(raw, max_pages), max_chars)
        return _clean(reader(raw), max_chars)
    except Exception as e:
        # A malformed or encrypted file is unreadable, not empty. Callers key
        # off the empty string and report rather than propose.
        log.warning("Extraction failed for %s: %s", url or "file", e)
        return ""
